# =============================================================================
# src/document_loaders/pptx_loader.py
# =============================================================================
"""
Loader para presentaciones PowerPoint (.pptx)
"""
from pathlib import Path
from typing import List, Dict, Any

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from .base_loader import BaseDocumentLoader, DocumentSection, ProcessedDocument


class PPTXLoader(BaseDocumentLoader):
    """Carga y procesa presentaciones PowerPoint"""

    def __init__(self):
        super().__init__()
        self.supported_extensions = {'.pptx'}

    def load(self, file_path: Path, original_filename: str = None) -> ProcessedDocument:
        """Carga una presentación PowerPoint"""
        prs = Presentation(file_path)

        # Extraer metadata
        metadata = self._extract_metadata(prs)

        # Extraer secciones (cada slide es una sección)
        sections = self.extract_sections(prs)

        # Generar contenido completo desde las secciones
        content = self._generate_full_content(sections)

        # Convertir a ruta relativa
        abs_path = file_path if file_path.is_absolute() else file_path.resolve()
        try:
            relative_path = abs_path.relative_to(Path.cwd())
        except ValueError:
            relative_path = abs_path

        return ProcessedDocument(
            file_path=str(relative_path),
            file_name=file_path.name,
                original_filename=original_filename or file_path.name,
            content=content,
            sections=sections,
            metadata=metadata
        )

    def _extract_metadata(self, prs: Presentation) -> Dict[str, Any]:
        """Extrae metadata de la presentación"""
        core_props = prs.core_properties

        metadata = {
            'author': core_props.author or '',
            'title': core_props.title or '',
            'subject': core_props.subject or '',
            'keywords': core_props.keywords or '',
            'comments': core_props.comments or '',
            'category': core_props.category or '',
            'created': core_props.created.isoformat() if core_props.created else '',
            'modified': core_props.modified.isoformat() if core_props.modified else '',
            'last_modified_by': core_props.last_modified_by or '',
            'slides_count': len(prs.slides),
            'slide_width': prs.slide_width,
            'slide_height': prs.slide_height
        }

        return metadata

    def _extract_text_from_shape(self, shape) -> str:
        """Extrae texto de una forma (shape) de PowerPoint"""
        text_content = []

        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                para_text = paragraph.text.strip()
                if para_text:
                    # Detectar nivel de indentación/jerarquía
                    level = paragraph.level
                    indent = '  ' * level
                    text_content.append(f"{indent}{para_text}")

        return '\n'.join(text_content) if text_content else ''

    def _extract_table_content(self, shape) -> str:
        """Extrae contenido de una tabla"""
        if not shape.has_table:
            return ''

        table = shape.table
        table_text = ['[TABLA]']

        for row in table.rows:
            row_text = ' | '.join(cell.text.strip() for cell in row.cells)
            if row_text.strip():
                table_text.append(row_text)

        table_text.append('[/TABLA]')
        return '\n'.join(table_text)

    def _extract_slide_notes(self, slide) -> str:
        """Extrae las notas del orador de una diapositiva"""
        if not slide.has_notes_slide:
            return ''

        notes_slide = slide.notes_slide
        notes_text_frame = notes_slide.notes_text_frame

        if notes_text_frame and notes_text_frame.text:
            return notes_text_frame.text.strip()

        return ''

    def extract_sections(self, prs: Presentation) -> List[DocumentSection]:
        """
        Extrae secciones de la presentación.
        Cada slide se considera una sección.
        """
        sections = []

        for slide_num, slide in enumerate(prs.slides, 1):
            slide_content = []
            slide_title = f"Slide {slide_num}"

            # Extraer contenido de todas las formas en orden
            for shape in slide.shapes:
                # Título del slide (generalmente el primer TextFrame)
                if shape.has_text_frame and hasattr(shape, 'text'):
                    text = shape.text.strip()
                    if text and len(slide_content) == 0:
                        # Primer texto, probablemente el título
                        slide_title = text if len(text) < 100 else f"Slide {slide_num}"
                        continue

                # Extraer texto de formas
                if shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX or \
                    shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER or \
                    shape.has_text_frame:
                    text = self._extract_text_from_shape(shape)
                    if text:
                        slide_content.append(text)

                # Extraer contenido de tablas
                elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                    table_content = self._extract_table_content(shape)
                    if table_content:
                        slide_content.append(table_content)

                # Extraer texto de grupos de formas
                elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                    group_text = self._extract_group_content(shape)
                    if group_text:
                        slide_content.append(group_text)

            # Extraer notas del orador
            notes = self._extract_slide_notes(slide)
            if notes:
                slide_content.append(f"\n[NOTAS DEL ORADOR]\n{notes}\n[/NOTAS]")

            # Crear la sección del slide
            if slide_content or slide_title != f"Slide {slide_num}":
                content_text = '\n\n'.join(slide_content).strip()

                sections.append(DocumentSection(
                    title=slide_title,
                    content=content_text if content_text else "[Slide sin contenido de texto]",
                    level=1,
                    metadata={
                        'slide_number': slide_num,
                        'has_notes': bool(notes)
                    }
                ))

        return sections

    def _extract_group_content(self, group_shape) -> str:
        """Extrae texto de un grupo de formas"""
        group_text = []

        for shape in group_shape.shapes:
            if shape.has_text_frame:
                text = self._extract_text_from_shape(shape)
                if text:
                    group_text.append(text)
            elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                # Recursivo para grupos anidados
                nested_text = self._extract_group_content(shape)
                if nested_text:
                    group_text.append(nested_text)

        return '\n'.join(group_text) if group_text else ''

    def _generate_full_content(self, sections: List[DocumentSection]) -> str:
        """Genera el contenido completo de la presentación desde las secciones"""
        full_content = []

        for section in sections:
            # Encabezado del slide
            slide_num = section.metadata.get('slide_number', '?')
            full_content.append(f"# Slide {slide_num}: {section.title}")

            # Contenido del slide
            if section.content:
                full_content.append(section.content)

            # Separador entre slides
            full_content.append("\n" + "=" * 80 + "\n")

        return '\n'.join(full_content).strip()
